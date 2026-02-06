"""
PowerPoint presentation (.pptx) text extraction using python-pptx.

Extracts:
- Slide titles
- Text from all shapes (text boxes, placeholders)
- Table content
- Speaker notes
"""

import io
import logging

from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_pptx_text(file_bytes: bytes) -> str:
    """
    Extract text content from a PowerPoint presentation (.pptx).

    Args:
        file_bytes: Raw PPTX file content

    Returns:
        Extracted text organized by slide.
        Returns empty string if extraction fails.
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = _extract_slide_text(slide, slide_num)
            if slide_text:
                text_parts.append(slide_text)

        full_text = "\n\n".join(text_parts)
        logger.info("Extracted %d characters from PPTX (%d slides)", len(full_text), len(prs.slides))
        return full_text

    except Exception as e:
        logger.error("Failed to extract PPTX: %s", e)
        return ""


def _extract_slide_text(slide, slide_num: int) -> str:
    """
    Extract text from a single slide.

    Includes title, body text, tables, and speaker notes.
    """
    lines: list[str] = []

    # Get slide title
    title = None
    if slide.shapes.title:
        title = slide.shapes.title.text.strip()
        if title:
            lines.append(f"Slide {slide_num}: {title}")
    else:
        lines.append(f"Slide {slide_num}")

    # Extract text from all shapes
    for shape in slide.shapes:
        # Skip the title shape (already processed)
        if shape == slide.shapes.title:
            continue

        # Text frames (text boxes, placeholders)
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text and text != title:  # Avoid duplicating title
                    lines.append(f"  {text}")

        # Tables
        if shape.has_table:
            table_text = _extract_table_text(shape.table)
            if table_text:
                lines.append(table_text)

    # Speaker notes
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            notes_text = notes_frame.text.strip()
            if notes_text:
                lines.append(f"  [Notes] {notes_text}")

    # Only return if we have content beyond just the slide number
    if len(lines) > 1 or (len(lines) == 1 and ":" in lines[0]):
        return "\n".join(lines)
    return ""


def _extract_table_text(table) -> str:
    """
    Convert a PowerPoint table to readable text.
    """
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        if any(cells):
            rows.append(cells)

    if not rows:
        return ""

    lines = ["  [Table]"]
    headers = rows[0] if rows else []

    for i, row in enumerate(rows[1:], start=1):
        pairs = []
        for j, cell in enumerate(row):
            if cell:
                header = headers[j] if j < len(headers) else f"Col{j+1}"
                pairs.append(f"{header}={cell}")
        if pairs:
            lines.append(f"    Row {i}: {', '.join(pairs)}")

    return "\n".join(lines) if len(lines) > 1 else ""
